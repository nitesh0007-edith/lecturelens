"""
Parsers for PDF, PPTX, Jupyter notebooks, and HTML → raw Document objects.
Each parser returns a list of RawPage objects preserving source structure.
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Literal


@dataclass
class RawPage:
    """A single logical unit from a source file before chunking."""

    text: str
    source_file: str
    page_or_slide: int
    doc_type: Literal["lecture", "lab", "past_paper", "cheatsheet", "coursework", "unknown"]
    title: str = ""
    extra_metadata: dict = field(default_factory=dict)


# ---------------------------------------------------------------------------
# PDF parser (PyMuPDF)
# ---------------------------------------------------------------------------

def parse_pdf(path: Path, doc_type: str = "lecture") -> list[RawPage]:
    """Extract pages from a PDF preserving heading structure."""
    try:
        import fitz  # PyMuPDF
    except ImportError as e:
        raise ImportError("pymupdf required: pip install pymupdf") from e

    pages: list[RawPage] = []
    doc = fitz.open(str(path))

    for page_num, page in enumerate(doc, start=1):
        blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
        lines: list[str] = []
        heading = ""

        for block in blocks:
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                spans = line.get("spans", [])
                if not spans:
                    continue
                text = " ".join(s["text"] for s in spans).strip()
                if not text:
                    continue
                # Heuristic: large font or bold → treat as heading
                max_size = max(s.get("size", 0) for s in spans)
                is_bold = any(s.get("flags", 0) & 16 for s in spans)
                if max_size >= 14 or is_bold:
                    if not heading:
                        heading = text
                lines.append(text)

        full_text = "\n".join(lines).strip()
        if full_text:
            pages.append(
                RawPage(
                    text=full_text,
                    source_file=str(path),
                    page_or_slide=page_num,
                    doc_type=doc_type,
                    title=heading,
                )
            )

    doc.close()
    return pages


# ---------------------------------------------------------------------------
# PPTX parser
# ---------------------------------------------------------------------------

def parse_pptx(path: Path, doc_type: str = "lecture") -> list[RawPage]:
    """One RawPage per slide — title + body + speaker notes."""
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ImportError("python-pptx required: pip install python-pptx") from e

    prs = Presentation(str(path))
    pages: list[RawPage] = []

    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts: list[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # Title placeholder
            if shape.shape_type == 13 or (hasattr(shape, "placeholder_format") and
                    shape.placeholder_format is not None and
                    shape.placeholder_format.idx == 0):
                title = text
            else:
                body_parts.append(text)

        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        parts = [p for p in [title] + body_parts + ([notes_text] if notes_text else []) if p]
        full_text = "\n".join(parts)

        if full_text.strip():
            pages.append(
                RawPage(
                    text=full_text,
                    source_file=str(path),
                    page_or_slide=idx,
                    doc_type=doc_type,
                    title=title,
                )
            )

    return pages


# ---------------------------------------------------------------------------
# Jupyter notebook parser
# ---------------------------------------------------------------------------

def parse_ipynb(path: Path, doc_type: str = "lab") -> list[RawPage]:
    """Pair markdown cell + following code cell(s) as one logical unit."""
    try:
        import nbformat
    except ImportError as e:
        raise ImportError("nbformat required: pip install nbformat") from e

    nb = nbformat.read(str(path), as_version=4)
    pages: list[RawPage] = []
    unit_idx = 0
    buffer: list[str] = []
    current_heading = ""

    def flush(idx: int, heading: str, buf: list[str]) -> RawPage | None:
        text = "\n".join(buf).strip()
        if not text:
            return None
        return RawPage(
            text=text,
            source_file=str(path),
            page_or_slide=idx,
            doc_type=doc_type,
            title=heading,
        )

    for cell in nb.cells:
        if cell.cell_type == "markdown":
            if buffer:
                p = flush(unit_idx, current_heading, buffer)
                if p:
                    pages.append(p)
                unit_idx += 1
                buffer = []
            src = cell.source.strip()
            # Extract heading from first # line
            heading_match = re.match(r"^#{1,3}\s+(.+)", src, re.MULTILINE)
            current_heading = heading_match.group(1) if heading_match else ""
            buffer.append(src)

        elif cell.cell_type == "code":
            src = cell.source.strip()
            if not src:
                continue
            # Strip large outputs; keep small ones
            output_text = ""
            for output in cell.outputs:
                text_out = ""
                if output.get("output_type") in ("stream", "execute_result", "display_data"):
                    text_out = "".join(
                        output.get("text", output.get("data", {}).get("text/plain", ""))
                    )
                if text_out and len(text_out) < 500:
                    output_text += f"\n# Output:\n{text_out}"
            buffer.append(f"```python\n{src}{output_text}\n```")

    if buffer:
        p = flush(unit_idx, current_heading, buffer)
        if p:
            pages.append(p)

    return pages


# ---------------------------------------------------------------------------
# HTML parser
# ---------------------------------------------------------------------------

def parse_html(path: Path, doc_type: str = "lecture") -> list[RawPage]:
    """Extract text from HTML lecture notes via BeautifulSoup."""
    try:
        from bs4 import BeautifulSoup
    except ImportError as e:
        raise ImportError("beautifulsoup4 required: pip install beautifulsoup4") from e

    html = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(html, "lxml")

    # Remove scripts, styles, nav
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()

    sections: list[RawPage] = []
    current_heading = ""
    current_parts: list[str] = []
    section_idx = 0

    def flush():
        nonlocal section_idx, current_heading, current_parts
        text = "\n".join(current_parts).strip()
        if text:
            sections.append(
                RawPage(
                    text=text,
                    source_file=str(path),
                    page_or_slide=section_idx,
                    doc_type=doc_type,
                    title=current_heading,
                )
            )
            section_idx += 1
        current_parts = []

    body = soup.find("body") or soup
    for tag in body.descendants:
        if not hasattr(tag, "name"):
            continue
        if tag.name in ("h1", "h2", "h3"):
            flush()
            current_heading = tag.get_text(" ", strip=True)
            current_parts.append(current_heading)
        elif tag.name in ("p", "li", "td", "th", "pre", "code"):
            text = tag.get_text(" ", strip=True)
            if text:
                current_parts.append(text)

    flush()

    if not sections:
        # Fallback: entire body as one page
        text = (soup.get_text("\n", strip=True) or "").strip()
        if text:
            sections.append(
                RawPage(text=text, source_file=str(path), page_or_slide=0, doc_type=doc_type)
            )

    return sections


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------

PAST_PAPER_SUFFIXES = {"past_paper", "exam", "pastpaper"}


def infer_doc_type(path: Path) -> str:
    name_lower = path.stem.lower()
    if any(kw in name_lower for kw in ("past_paper", "exam", "pastpaper", "2015", "2016",
                                        "2017", "2018", "2019", "2020", "2021", "2022",
                                        "2023", "2024", "2025")):
        return "past_paper"
    if any(kw in name_lower for kw in ("lab", "tutorial", "exercise", "worksheet")):
        return "lab"
    if any(kw in name_lower for kw in ("cheatsheet", "cheat_sheet", "summary", "formula")):
        return "cheatsheet"
    if any(kw in name_lower for kw in ("coursework", "assignment", "cw")):
        return "coursework"
    return "lecture"


def parse_file(path: Path) -> list[RawPage]:
    """Auto-dispatch to the right parser based on file extension."""
    doc_type = infer_doc_type(path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return parse_pdf(path, doc_type)
    elif suffix in (".pptx", ".ppt"):
        return parse_pptx(path, doc_type)
    elif suffix == ".ipynb":
        return parse_ipynb(path, doc_type)
    elif suffix in (".html", ".htm"):
        return parse_html(path, doc_type)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")
